"""
Generate BioReasonCheck-FI pitch deck as a PowerPoint file.
Usage: python scripts/make_ppt.py
Output: outputs/BioReasonCheck_FI.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Palette ──────────────────────────────────────────────────────────────────
DARK_BG    = RGBColor(0x0D, 0x1B, 0x2A)   # near-black navy
ACCENT     = RGBColor(0x00, 0xC8, 0xFF)   # cyan
ACCENT2    = RGBColor(0xFF, 0x6B, 0x35)   # orange (highlight numbers)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xB0, 0xC4, 0xD8)
CARD_BG    = RGBColor(0x16, 0x2A, 0x3E)   # slightly lighter panel

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]   # blank
    slide = prs.slides.add_slide(blank_layout)
    # Dark background rectangle
    bg = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        0, 0, SLIDE_W, SLIDE_H
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()
    return slide


def add_textbox(slide, left, top, width, height,
                text, size, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_para(tf, text, size, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, space_before=0, italic=False):
    """Add a paragraph to an existing text frame."""
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def accent_bar(slide, y=Inches(0.55)):
    """Thin cyan horizontal rule under the title area."""
    bar = slide.shapes.add_shape(1, Inches(0.4), y, Inches(12.53), Pt(2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


def slide_number(slide, n: int):
    add_textbox(slide,
                Inches(12.5), Inches(7.1), Inches(0.7), Inches(0.3),
                str(n), 9, color=LIGHT_GREY, align=PP_ALIGN.RIGHT)


def header(slide, title: str, subtitle: str = ""):
    add_textbox(slide,
                Inches(0.4), Inches(0.12), Inches(12), Inches(0.5),
                title, 26, bold=True, color=ACCENT)
    if subtitle:
        add_textbox(slide,
                    Inches(0.4), Inches(0.55), Inches(12), Inches(0.35),
                    subtitle, 14, color=LIGHT_GREY)


def card(slide, left, top, width, height):
    """Semi-transparent card panel."""
    sh = slide.shapes.add_shape(1, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = CARD_BG
    sh.line.color.rgb = ACCENT
    sh.line.width = Pt(0.75)
    return sh


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — TITLE
# ─────────────────────────────────────────────────────────────────────────────
def slide_title(prs):
    slide = blank_slide(prs)

    # Centered project name
    add_textbox(slide,
                Inches(1), Inches(1.8), Inches(11.33), Inches(1.1),
                "BioReasonCheck-FI", 48, bold=True,
                color=ACCENT, align=PP_ALIGN.CENTER)

    add_textbox(slide,
                Inches(1), Inches(3.0), Inches(11.33), Inches(0.6),
                "Format Instability Benchmark for LongevityLLM", 22,
                color=WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide,
                Inches(1), Inches(3.7), Inches(11.33), Inches(0.45),
                "Track 01 · Insilico Medicine Hackathon 2026", 16,
                color=LIGHT_GREY, align=PP_ALIGN.CENTER, italic=True)

    # Stat pill row
    stats = [
        ("95%", "Format Instability"),
        ("270", "Benchmark Prompts"),
        ("60", "Unique Genes"),
        ("p=0.021", "McNemar Test"),
    ]
    pill_w = Inches(2.8)
    pill_h = Inches(0.9)
    gap    = Inches(0.2)
    start  = Inches(0.76)
    y      = Inches(5.0)

    for i, (num, label) in enumerate(stats):
        x = start + i * (pill_w + gap)
        c = card(slide, x, y, pill_w, pill_h)
        add_textbox(slide, x, y + Inches(0.05), pill_w, Inches(0.45),
                    num, 24, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)
        add_textbox(slide, x, y + Inches(0.48), pill_w, Inches(0.35),
                    label, 11, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

    add_textbox(slide,
                Inches(1), Inches(6.8), Inches(11.33), Inches(0.3),
                "github.com/aakashak2000/BioReasonCheck", 11,
                color=LIGHT_GREY, align=PP_ALIGN.CENTER)

    slide_number(slide, 1)
    return slide


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — THE PROBLEM
# ─────────────────────────────────────────────────────────────────────────────
def slide_problem(prs):
    slide = blank_slide(prs)
    header(slide, "The Problem", "Same gene. Same biology. Different answer.")
    accent_bar(slide, y=Inches(0.88))
    slide_number(slide, 2)

    # Big quote / analogy
    card(slide, Inches(0.4), Inches(1.1), Inches(12.53), Inches(1.55))
    add_textbox(slide,
                Inches(0.6), Inches(1.18), Inches(12.1), Inches(1.4),
                '"Imagine asking a specialist: Does this patient have condition X? '
                'They say NO. You show them a form with the same question as multiple choice. '
                'They circle YES. Same specialist. Same patient. Same biology. '
                'That\'s what we found with LongevityLLM — in 19 of 20 test cases."',
                14, italic=True, color=WHITE)

    # Three columns
    col_w = Inches(3.9)
    col_h = Inches(3.2)
    col_y = Inches(2.9)
    cols = [
        ("Binary API call", 'Gene X\nsenescent?', "NOT_SUPPORTED ✗", ACCENT2),
        ("Ternary classifier", 'Gene X\nsenescent?', "INSUFFICIENT_EVIDENCE ✗", LIGHT_GREY),
        ("MCQ interface", 'Gene X\nsenescent?', "SUPPORTED ✓  (but first two wrong)", ACCENT),
    ]
    for i, (title, q, ans, col) in enumerate(cols):
        x = Inches(0.4) + i * (col_w + Inches(0.17))
        card(slide, x, col_y, col_w, col_h)
        add_textbox(slide, x + Inches(0.1), col_y + Inches(0.1), col_w - Inches(0.2), Inches(0.4),
                    title, 13, bold=True, color=col)
        add_textbox(slide, x + Inches(0.1), col_y + Inches(0.6), col_w - Inches(0.2), Inches(0.8),
                    q, 18, color=LIGHT_GREY)
        add_textbox(slide, x + Inches(0.1), col_y + Inches(1.5), col_w - Inches(0.2), Inches(1.5),
                    ans, 15, bold=True, color=col)

    add_textbox(slide,
                Inches(0.4), Inches(6.2), Inches(12.5), Inches(0.4),
                "No prior evaluation had measured this.  We built the test.", 14,
                bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    return slide


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — WHAT WE BUILT
# ─────────────────────────────────────────────────────────────────────────────
def slide_benchmark(prs):
    slide = blank_slide(prs)
    header(slide, "What We Built",
           "270 prompts · 60 genes · 3 formats · 5 categories · pre-registered test split")
    accent_bar(slide, y=Inches(0.88))
    slide_number(slide, 3)

    # Left: data design
    left_x = Inches(0.4)
    left_w = Inches(5.8)
    card(slide, left_x, Inches(1.1), left_w, Inches(5.8))

    add_textbox(slide, left_x + Inches(0.15), Inches(1.2), left_w - Inches(0.3), Inches(0.4),
                "Data Design", 15, bold=True, color=ACCENT)

    rows = [
        ("5 Gene Categories", "CellAge senescence, OpenGenes lifespan, mixed-evidence, metabolic hard negatives, invented hallucination traps"),
        ("Ground Truth", "CellAge & OpenGenes database exports — not PubMed, not Wikipedia"),
        ("3 Claim Tiers", "Database-membership · False-strong (requires reasoning) · Hallucination traps (zero training signal)"),
        ("Pre-registered Split", "Deterministic gene hash (MD5 % 10 ≥ 7 → test). Assigned before any analysis ran."),
    ]
    y_off = Inches(1.7)
    for label, desc in rows:
        add_textbox(slide, left_x + Inches(0.15), y_off, left_w - Inches(0.3), Inches(0.25),
                    label, 12, bold=True, color=ACCENT2)
        add_textbox(slide, left_x + Inches(0.15), y_off + Inches(0.26), left_w - Inches(0.3), Inches(0.5),
                    desc, 11, color=LIGHT_GREY)
        y_off += Inches(0.88)

    # Right: format diversity
    right_x = Inches(6.53)
    right_w = Inches(6.4)
    card(slide, right_x, Inches(1.1), right_w, Inches(5.8))

    add_textbox(slide, right_x + Inches(0.15), Inches(1.2), right_w - Inches(0.3), Inches(0.4),
                "Format Diversity (the rubric axis)", 15, bold=True, color=ACCENT)

    formats = [
        ("Binary", "YES / NO", "Simple existence check", ACCENT2),
        ("Ternary", "SUPPORTED / NOT_SUPPORTED / INSUFFICIENT_EVIDENCE", "3-class reasoning label", ACCENT),
        ("MCQ", "4 options, correct position rotated", "Positional heuristic test", LIGHT_GREY),
    ]
    y_off = Inches(1.75)
    for fmt, labels, note, col in formats:
        card(slide, right_x + Inches(0.15), y_off, right_w - Inches(0.3), Inches(1.35))
        add_textbox(slide, right_x + Inches(0.3), y_off + Inches(0.05), Inches(1.2), Inches(0.35),
                    fmt, 16, bold=True, color=col)
        add_textbox(slide, right_x + Inches(0.3), y_off + Inches(0.42), right_w - Inches(0.6), Inches(0.4),
                    labels, 11, color=WHITE)
        add_textbox(slide, right_x + Inches(0.3), y_off + Inches(0.85), right_w - Inches(0.6), Inches(0.35),
                    note, 10, italic=True, color=LIGHT_GREY)
        y_off += Inches(1.52)

    return slide


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — KEY FINDINGS
# ─────────────────────────────────────────────────────────────────────────────
def slide_findings(prs):
    slide = blank_slide(prs)
    header(slide, "Key Findings", "Format instability is statistically significant and mechanistically explained")
    accent_bar(slide, y=Inches(0.88))
    slide_number(slide, 4)

    metrics = [
        ("95.0%", "Format Instability Rate\n[Bootstrap CI: 85%–100%]", ACCENT2),
        ("−0.328", "Cohen's κ  binary vs MCQ\n(worse than chance)", ACCENT),
        ("0.760", "Cramér's V\nformat → label association", ACCENT2),
        ("p = 0.021", "McNemar test\n(all 60 genes)", ACCENT),
        ("84.6%\nvs 34.6%", "Claude Sonnet 4.6 MCQ\nvs L-LLM MCQ", ACCENT2),
        ("72.7%", "Consistency mismatch\nerror rate (22 cases)", ACCENT),
    ]

    card_w = Inches(3.95)
    card_h = Inches(2.5)
    gap    = Inches(0.15)

    for i, (num, label, col) in enumerate(metrics):
        row = i // 3
        col_i = i % 3
        x = Inches(0.4) + col_i * (card_w + gap)
        y = Inches(1.1) + row * (card_h + gap)
        card(slide, x, y, card_w, card_h)
        add_textbox(slide, x + Inches(0.1), y + Inches(0.15), card_w - Inches(0.2), Inches(1.0),
                    num, 28, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.1), y + Inches(1.25), card_w - Inches(0.2), Inches(1.1),
                    label, 12, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

    return slide


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — THE MECHANISM
# ─────────────────────────────────────────────────────────────────────────────
def slide_mechanism(prs):
    slide = blank_slide(prs)
    header(slide, "The Mechanism", "Format decides the answer. The biology is irrelevant.")
    accent_bar(slide, y=Inches(0.88))
    slide_number(slide, 5)

    # Binary bias panel
    card(slide, Inches(0.4), Inches(1.1), Inches(5.9), Inches(4.8))
    add_textbox(slide, Inches(0.55), Inches(1.2), Inches(5.6), Inches(0.4),
                "Binary Format Bias", 15, bold=True, color=ACCENT2)
    add_textbox(slide, Inches(0.55), Inches(1.72), Inches(5.6), Inches(0.8),
                "92%", 52, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.55), Inches(2.6), Inches(5.6), Inches(0.4),
                "of answers are NOT_SUPPORTED", 13, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.55), Inches(3.1), Inches(5.6), Inches(0.4),
                "False Negative Rate: 83%", 13, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.55), Inches(3.65), Inches(5.6), Inches(1.1),
                'Binary triggers a "skeptical default" — when in doubt, say no.\nThe gene\'s biology plays no role.', 12,
                color=LIGHT_GREY, align=PP_ALIGN.CENTER, italic=True)

    # MCQ bias panel
    card(slide, Inches(6.63), Inches(1.1), Inches(5.9), Inches(4.8))
    add_textbox(slide, Inches(6.78), Inches(1.2), Inches(5.6), Inches(0.4),
                "MCQ Format Bias", 15, bold=True, color=ACCENT)
    add_textbox(slide, Inches(6.78), Inches(1.72), Inches(5.6), Inches(0.8),
                "81%", 52, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(6.78), Inches(2.6), Inches(5.6), Inches(0.4),
                "of answers are SUPPORTED", 13, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(6.78), Inches(3.1), Inches(5.6), Inches(0.4),
                "False Positive Rate: 86%", 13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(6.78), Inches(3.65), Inches(5.6), Inches(1.1),
                'MCQ triggers an "optimistic default" — pick the positive-sounding option.\nCompletely opposite heuristic.', 12,
                color=LIGHT_GREY, align=PP_ALIGN.CENTER, italic=True)

    # Bottom bar
    card(slide, Inches(0.4), Inches(6.05), Inches(12.1), Inches(0.8))
    add_textbox(slide, Inches(0.55), Inches(6.1), Inches(12.0), Inches(0.65),
                "Cramér's V = 0.760  ·  Format alone predicts label choice  ·  Cohen's κ (binary vs MCQ) = −0.328 — worse than chance agreement",
                13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    return slide


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — SCORER & RECOMMENDATIONS
# ─────────────────────────────────────────────────────────────────────────────
def slide_scorer_recs(prs):
    slide = blank_slide(prs)
    header(slide, "Reasoning Scorer + Recommendations",
           "Automatic error detection from chain-of-thought · Four actionable fixes")
    accent_bar(slide, y=Inches(0.88))
    slide_number(slide, 6)

    # Left: scorer
    lx = Inches(0.4)
    lw = Inches(6.1)
    card(slide, lx, Inches(1.1), lw, Inches(5.8))
    add_textbox(slide, lx + Inches(0.15), Inches(1.2), lw - Inches(0.3), Inches(0.4),
                "The Reasoning Scorer", 15, bold=True, color=ACCENT)
    add_textbox(slide, lx + Inches(0.15), Inches(1.65), lw - Inches(0.3), Inches(0.4),
                "No human labels needed — reads the model's own chain-of-thought", 11,
                color=LIGHT_GREY, italic=True)

    scorer_items = [
        ("78%", "of ternary traces use SENESCENCE as a gene symbol\n→ process-as-gene confusion"),
        ("72.7%", "error rate when reasoning contradicts output\n(22 consistency mismatch cases)"),
        ("85%", "scorer recall  ·  F1 = 0.694\ndeployable as runtime filter today"),
        ("Free DPO pairs", "each flagged trace = preference pair\nfor future fine-tuning"),
    ]
    y_off = Inches(2.15)
    for num, desc in scorer_items:
        add_textbox(slide, lx + Inches(0.15), y_off, Inches(1.5), Inches(0.35),
                    num, 15, bold=True, color=ACCENT2)
        add_textbox(slide, lx + Inches(1.7), y_off, lw - Inches(1.9), Inches(0.7),
                    desc, 11, color=WHITE)
        y_off += Inches(0.95)

    # Right: recommendations
    rx = Inches(6.83)
    rw = Inches(6.1)
    card(slide, rx, Inches(1.1), rw, Inches(5.8))
    add_textbox(slide, rx + Inches(0.15), Inches(1.2), rw - Inches(0.3), Inches(0.4),
                "Recommendations for Insilico", 15, bold=True, color=ACCENT)

    recs = [
        ("1", "Avoid binary format", "92% NOT_SUPPORTED bias regardless of true label — systematically unreliable"),
        ("2", "Rotate MCQ positions", "Option D was never selected correctly (0/5 cases) — positional heuristics are strong"),
        ("3", "Deploy scorer as runtime filter", "Flags 85% of errors automatically — no gold labels, no annotation pipeline needed"),
        ("4", "Add FIR to fine-tuning eval", "Format stability must be measured alongside accuracy in every future L-LLM training run"),
    ]
    y_off = Inches(1.75)
    for num, title, desc in recs:
        add_textbox(slide, rx + Inches(0.15), y_off, Inches(0.45), Inches(0.45),
                    num, 20, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)
        add_textbox(slide, rx + Inches(0.7), y_off, rw - Inches(0.9), Inches(0.3),
                    title, 13, bold=True, color=WHITE)
        add_textbox(slide, rx + Inches(0.7), y_off + Inches(0.32), rw - Inches(0.9), Inches(0.55),
                    desc, 11, color=LIGHT_GREY)
        y_off += Inches(1.1)

    return slide


# ─────────────────────────────────────────────────────────────────────────────
# MAIN
# ─────────────────────────────────────────────────────────────────────────────
def main():
    out_path = os.path.join(os.path.dirname(__file__), "..", "outputs", "BioReasonCheck_FI.pptx")
    prs = new_prs()

    slide_title(prs)
    slide_problem(prs)
    slide_benchmark(prs)
    slide_findings(prs)
    slide_mechanism(prs)
    slide_scorer_recs(prs)

    prs.save(out_path)
    print(f"Saved: {os.path.abspath(out_path)}")


if __name__ == "__main__":
    main()
